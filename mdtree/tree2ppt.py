import datetime
import json
import os
import random
from enum import Enum
from io import BytesIO
import cv2
import urllib
import numpy as np
import requests as req
from PIL import Image
from io import BytesIO
import markdown
from PIL.ImageQt import rgb
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt,Cm
from pptx.dml.color import RGBColor
from pptx.text.text import Font

from mdtree.parser import parse_string, Out, Heading
from mdtree.utils import get_random_theme, get_random_file, read_md_file
import os
from mdtree.img_search import get_img
import os
class Tree2PPT:
    prs: Presentation = None
    md_str: str = None
    out: Out = None
    tree: Heading = None
    theme: str = None

    def __init__(self, md_str1,theme_path,save_path=""):
        self.theme = theme_path
        self.theme_param_path = os.path.join(self.theme,"mode.json")
        with open(self.theme_param_path,encoding="utf-8") as f:
            self.theme_param = json.load(f)

        print(self.theme)
        self.init_pptx(theme_path)
        self.init_markdown(md_str1)
        self.ppt_main_theme = self.tree.text
        THEME_MD2Slide(self.prs, self.theme, self.ppt_main_theme)
        keywords = self.ppt_main_theme
        if self.theme_param["main_page"].get("img_info"):
            try:
                self.img_dicts = get_img(keywords)
                # self.img_dicts = ddg_images(keywords, region='wt-wt', safesearch='Off', size=None,color='Monochrome', type_image=None, layout=None, license_image=None, max_results=300)
                print("一共检索到的图片数量为",len(self.img_dicts))
            except:
                self.img_dicts=[]
        else:
            self.img_dicts=[]
        # {'height': 1000, 'image': 'h', 'source': 'Bing', 'thumbnail': 'httpApi', 'title': 'Liberty Tree 1765 Nthe Large Elm Tree At Boylston Market ...', 'url': 'htt650000', 'width': 1000},
        self.traverse_tree(self.tree)
        now = datetime.datetime.now().timestamp()
        if not os.path.exists('./myppt'):
            os.makedirs('./myppt')
        if save_path=="":
            self.save_path = './myppt/test.pptx'
        else:
            self.save_path = save_path

        self.prs.save(self.save_path)
        pass

    def init_pptx(self,theme_path = "../my_ppt_mode/1"):
        prs = Presentation()
        prs.slide_height = Cm(self.theme_param["slide_size"]["height"])
        prs.slide_width = Cm(self.theme_param["slide_size"]["width"])
        self.theme = theme_path
        self.prs = prs

    def init_markdown(self, md_str):
        self.md_str = md_str
        self.out = parse_string(md_str)
        self.tree = self.out.main

    def traverse_tree(self, heading):
        if heading is not None and (heading.source is None or heading.source.strip() == ''):#当主题下面是分主题时
            content = ""
            if len(heading.children)>0:
                for child in heading.children:
                    content = content + child.text + "\n"

            content = content.replace("引文","").replace("总结","").strip()
            if heading.text_source[:2]=="# ":
                MD2Slide(self.prs, self.theme, "目录",content=content)
            else:
                MD2Slide(self.prs, self.theme, heading.text,content=content)
        elif heading is not None:
            content_list = heading.source.split("\n\n")
            content_max_word_num = 250
            begin_index = 0
            input_text_list = []
            while begin_index < len(content_list):
                # print("begin_index", begin_index, "所有元素个数", len(ele_text_list))
                input_token_i = 0
                input_text = ""
                for i in range(begin_index, len(content_list)):
                    if len(content_list[i]) + input_token_i < content_max_word_num:
                        input_text += "\n"+content_list[i]
                        input_token_i += len(content_list[i])
                        begin_index = i
                    else:
                        if input_token_i==0:
                            input_text +=content_list[i]
                        break
                begin_index += 1
                input_text_list.append(input_text)

            for content_i in range(len(input_text_list)):
                if content_i==0:
                    if self.theme_param["main_page"].get("img_info") and len(self.img_dicts)>0:
                        img_dict = self.img_dicts[random.randint(0, len(self.img_dicts) - 1)]
                    else:
                        img_dict = {}
                    MD2Slide(self.prs, self.theme, heading.text, img_dict=img_dict,content=input_text_list[content_i].strip())
                else:
                    if self.theme_param["main_page"].get("img_info") and len(self.img_dicts)>0:
                        img_dict = self.img_dicts[random.randint(0, len(self.img_dicts) - 1)]
                    else:
                        img_dict = {}
                    MD2Slide(self.prs, self.theme, "", img_dict=img_dict,content=input_text_list[content_i].strip())
            # MD2Slide(self.prs, self.theme, heading.text, content=heading.source)
        else:
            return

        # self.make_slide_demo(self.prs, heading.text, heading.source)
        if heading.children is not []:
            for child in heading.children:
                self.traverse_tree(child)

    def save_stream(self):
        stream = BytesIO()
        self.prs.save(stream)
        stream.seek(0)  # Reset the stream position to the beginning
        return stream


class MarkdownCategory:
    TITLE = "#"
    CONTENT = "<p>"

    pass


class MD2Slide:
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "黑体"
    font_title_size: Pt = Pt(26)
    font_content_size: Pt = Pt(14)
    font_title_color: rgb = RGBColor(0, 0, 0)
    font_content_color: rgb = RGBColor(0, 0, 0)
    title_box = (Inches(0.3), Inches(0.3), Cm(24.24), Inches(0.8))
    content_box = (Cm(2.54), Cm(4.12),Cm(20.32), Cm(12.70))

    def __init__(self, presentation, theme_path, title, content, *args,img_dict={}, **kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.title = title
        self.content = content
        self.theme = theme_path
        self.theme_param_path = os.path.join(self.theme,"mode.json")
        with open(self.theme_param_path,encoding="utf-8") as f:
            self.theme_param = json.load(f)
        print(self.theme)
        page_params = self.theme_param["main_page"]

        self.title_box = (Cm(page_params["title_info"]["pos_x"]),Cm(page_params["title_info"]["pos_y"]),Cm(page_params["title_info"]["width"]),Cm(page_params["title_info"]["height"]))
        self.content_box = (Cm(page_params["content_info"]["pos_x"]),Cm(page_params["content_info"]["pos_y"]),Cm(page_params["content_info"]["width"]),Cm(page_params["content_info"]["height"]))
        self.font_content_size = Pt(page_params["content_info"]["font_size"])
        self.font_title_size = Pt(page_params["title_info"]["font_size"])
        self.font_name = page_params["title_info"]["font_name"]
        if self.title =="目录":
            page_params = self.theme_param["catalog_page"]
            self.title_box = (Cm(page_params["title_info"]["pos_x"]),Cm(page_params["title_info"]["pos_y"]),Cm(page_params["title_info"]["width"]),Cm(page_params["title_info"]["height"]))
            self.font_title_size = Pt(page_params["title_info"]["font_size"])
            self.font_name = page_params["title_info"]["font_name"]

        self.img_theme= theme_path+"/"+"img"
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()
        if page_params.get("img_info") and img_dict:
            self.img_url = img_dict["thumbnail"]
            # img_orginal_h = img_dict["height"]
            # img_orginal_w = img_dict["width"]
            # img_h = float(page_params["img_info"]["width"])*float(img_orginal_h)/float(img_orginal_w)
            self.img_box = (Cm(page_params["img_info"]["pos_x"]),Cm(page_params["img_info"]["pos_y"]),Cm(page_params["img_info"]["width"]),Cm(page_params["img_info"]["height"]))
            self.init_img()

    def init_slide(self):
        # placeholder1 = self.slide.placeholders[1]
        path = get_random_file(self.img_theme)
        img_box = (Cm(0),Cm(0),self.presentation.slide_width,self.presentation.slide_height)
        picture = self.slide.shapes.add_picture(path,*img_box)
        # placeholder2 = self.slide.placeholders[2]
        # placeholder2.element.getparent().remove(placeholder2.element)
        # 2、设置占位符宽高
        # picture.left = 0
        # picture.top = 0
        # picture.width = self.presentation.slide_width
        # picture.height = self.presentation.slide_height

    def init_img(self):
        # encoding:utf-8

        response = req.get(self.img_url)
        image = Image.open(BytesIO(response.content))
        image.save("temp_img.png")
        picture = self.slide.shapes.add_picture("temp_img.png",*self.img_box)



    def init_font(self, **kwargs):
        if 'font_name' in kwargs:
            self.font_name = kwargs['font_name']
        if 'font_title_size' in kwargs:
            self.font_title_size = kwargs['font_title_size']
        if 'font_content_size' in kwargs:
            self.font_content_size = kwargs['font_content_size']
        if 'font_title_color' in kwargs:
            self.font_title_color = kwargs['font_title_color']
        if 'font_content_color' in kwargs:
            self.font_content_color = kwargs['font_content_color']
        if 'content_box' in kwargs:
            self.content_box = kwargs['content_box']
        if 'title_box' in kwargs:
            self.title_box = kwargs['title_box']

    def get_font(self, font: Font, category: str):

        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.bold = True
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        # if self.title not in ["首段","结尾"]:
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(*self.title_box)
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # 添加标题
        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def init_content(self):
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(*self.content_box)
        tf = text_box_content.text_frame
        tf.clear()  # Clear existing content
        # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True
        # 添加正文
        paragraph = tf.paragraphs[0]
        self.content = self.content.replace("<p>", "").replace("</p>", "\n")
        self.content = self.content.replace("\n\n","\n").replace("\n","\n\n")
        paragraph.text = self.content
        self.processing_md_str(self.content.replace("<p>", "").replace("</p>", "\n"))
        # TODO 处理正文
        self.get_font(paragraph.font, MarkdownCategory.CONTENT)
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def processing_md_str(self, md_str):
        print(md_str)
        md = markdown.Markdown()
        html1 = md.convert(md_str)
        print(html1)






class THEME_MD2Slide:
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "黑体"
    font_title_color: rgb = RGBColor(0, 0, 0)
    font_title_size:Pt = Pt(40)
    title_box = (Cm(2.81), Cm(5.44), Cm(21.59), Cm(4.08))

    def __init__(self, presentation, theme_path, title, *args, **kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.title = title
        self.theme = theme_path
        self.theme_param_path = os.path.join(self.theme,"mode.json")
        with open(self.theme_param_path,encoding="utf-8") as f:
            self.theme_param = json.load(f)
        print(self.theme)
        first_page_params = self.theme_param["first_page"]
        self.title_box = (Cm(first_page_params["title_info"]["pos_x"]),Cm(first_page_params["title_info"]["pos_y"]),Cm(first_page_params["title_info"]["width"]),Cm(first_page_params["title_info"]["height"]))
        self.font_title_size = Pt(first_page_params["title_info"]["font_size"])
        self.font_name = first_page_params["title_info"]["font_name"]
        self.img_theme= theme_path+"/"+"img"
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()

    def init_slide(self):
        # placeholder1 = self.slide.placeholders[1]
        if os.path.exists(os.path.join(self.theme,"title.jpg")):
            path = os.path.join(self.theme,"title.jpg")
        elif os.path.exists(os.path.join(self.theme,"title.png")):
            path = os.path.join(self.theme,"title.png")
        else:
            path = get_random_file(self.img_theme)
        left, top, width, height = (Cm(0),Cm(0),Cm(25.4),Cm(14.29))
        picture = self.slide.shapes.add_picture(path,left, top, width, height)
        # picture = placeholder1.insert_picture(path)
        # placeholder2 = self.slide.placeholders[2]
        # placeholder2.element.getparent().remove(placeholder2.element)
        # 2、设置占位符宽高
        picture.left = 0
        picture.top = 0
        picture.width = self.presentation.slide_width
        picture.height = self.presentation.slide_height

    def init_font(self, **kwargs):
        if 'font_name' in kwargs:
            self.font_name = kwargs['font_name']
        if 'font_title_size' in kwargs:
            self.font_title_size = kwargs['font_title_size']
        if 'font_content_size' in kwargs:
            self.font_content_size = kwargs['font_content_size']
        if 'font_title_color' in kwargs:
            self.font_title_color = kwargs['font_title_color']

    def get_font(self, font: Font, category: str):
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
            font.bold = True
    def init_title(self):
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(*self.title_box)
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        # tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True #允许自动换行

        # 添加标题
        paragraph = tf.paragraphs[0]
        # paragraph.alignment = PP_ALIGN.CENTER
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

